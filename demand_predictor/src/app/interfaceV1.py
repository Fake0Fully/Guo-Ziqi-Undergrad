import os
import sys
from kivy.app import App
from kivy.uix.filechooser import FileChooserListView
from kivy.uix.floatlayout import FloatLayout
from kivy.uix.gridlayout import GridLayout
from kivy.uix.boxlayout import BoxLayout
from kivy.uix.dropdown import DropDown
from kivy.uix.togglebutton import ToggleButton
from kivy.uix.textinput import TextInput
from kivy.uix.popup import Popup
from kivy.graphics import Rectangle
from kivy.uix.scrollview import ScrollView
from kivy.uix.label import Label
from kivy.uix.button import Button
from kivy.lang import Builder
from pptx import Presentation
from pptx.util import Inches, Pt
from read.macro_utils import MacroHelper
import prediction.predictor
import prediction.utils
import read.data_parser
import matplotlib.pyplot as plt
import pandas as pd
import numpy as np
import csv
import datetime
import copy

# TODO: add capacity * 13 to batch, handle exceptions (no data), upload file
_file = os.path.abspath(sys.argv[0])
PATH = os.path.dirname(_file)
Builder.load_file("interface.kv")

REF_PATH = {"ATV_All_All_All": PATH[:-8] + "/sample_data/atv.csv",
            "PMM_All_All_All": PATH[:-8] + "/sample_data/pmm.csv",
            "CCS_All_All_All": PATH[:-8] + "/sample_data/ccs.csv",
            "IPC_All_All_All": PATH[:-8] + "/sample_data/ipc.csv",
            "ATV_All_LQFP_All": PATH[:-8] + "/sample_data/atv_lqfp.csv",
            "ATV_All_TDSON_All": PATH[:-8] + "/sample_data/atv_tdson.csv",
            "ATV_All_TO252_All": PATH[:-8] + "/sample_data/atv_to252.csv",
            "ATV_All_VQFN_All": PATH[:-8] + "/sample_data/atv_vqfn.csv",
            "PMM_All_All_TSDSON": PATH[:-8] + "/sample_data/pmm_tsdson.csv",
            "PMM_All_VQFN_All": PATH[:-8] + "/sample_data/pmm_vqfn.csv",
            "PMM_All_All_VSON": PATH[:-8] + "/sample_data/pmm_vson.csv",
            "PMM_All_WLL_All": PATH[:-8] + "/sample_data/pmm_wll.csv"
            }
MAPPING_PATH = PATH[:-8] + "/sample_data/Multiple Factor Filtering.xlsx"
DATA_FILE = PATH[:-8] + "/sample_data/data.xlsm"
CONFIG_FILE = PATH[:-8] + "\\config.txt"
PIVOT_FIELDS = ["Division_short", "PL", "Agg_Package_PL", "PackageClass_PL"]

download_identifiers = []
download_buttons = []
download_details = []
download_stats = []
display_stats = []
download_images = []
display_images = []
window_counter = -1


class UI(FloatLayout):
    def __init__(self, **kwargs):
        super(UI, self).__init__(**kwargs)

        self.division_list = ["ATV", "PMM", "CCS", "IPC"]
        self.pkc_list = None
        self.pkca_list = None
        self.pl_list = None

        self.macro_helper = MacroHelper(PATH)
        self.mapping = None

        self.division_drop_down = None
        self.pl_drop_down = None
        self.aggPkc_drop_down = None
        self.pkgClass_drop_down = None

        self.scrollable_label = ScrollableLabel()
        self.ids.documentCounter.text = str(len(download_identifiers))

        self.ratio = None
        self.identifier = None
        self.results = None
        self.upload_popup = UploadPopUp()

    def reset_filters(self):
        success = self.macro_helper.execute_macro("change_pivot", "Division_short", ["All"])
        self.macro_helper.execute_macro("change_pivot", "PL", ["All"])
        self.macro_helper.execute_macro("change_pivot", "Agg_Package_PL", ["All"])
        self.macro_helper.execute_macro("change_pivot", "PackageClass_PL", ["All"])

        if success == "NA":
            print ("No workbooks opened")
            popup = Popup(title="Status", content=Label(text='No workbooks opened'), size_hint=(None, None),
                          size=(200, 100))
            popup.open()

    def init_filters(self):

        self.reset_filters()
        df = pd.ExcelFile(MAPPING_PATH)
        self.mapping = df.parse("mapping", index_col=None)

        self.macro_helper.execute_macro("get_filter", "PL")
        self.pl_list = open(PATH + "\\read\\" + "PL_output.txt", "r").readline().decode("utf-16").split("@")[:-1]
        self.pl_list = [x.encode("UTF-8") for x in self.pl_list]
        self.pl_list.sort()
        self.pkc_list = list(set([x.encode("UTF-8") for x in list(self.mapping["PKC"])]))
        self.pkc_list.sort()
        self.pkca_list = list(set([x.encode("UTF-8") for x in list(self.mapping["PKCA"])]))
        self.pkca_list.sort()

        self.ids.divisionButton.text = "Division (" + str(len(self.division_list)) + ")"
        self.ids.plButton.text = "PL (" + str(len(self.pl_list)) + ")"
        self.ids.pkcaButton.text = "PKCA (" + str(len(self.pkca_list)) + ")"
        self.ids.pkcButton.text = "PKC (" + str(len(self.pkc_list)) + ")"

    def close_workbook(self):
        self.macro_helper.terminate()
        print ("Workbook Closed")
        user_interface.ui.ids.divisionButton.disabled = True
        user_interface.ui.ids.pkcaButton.disabled = True
        user_interface.ui.ids.pkcButton.disabled = True
        user_interface.ui.ids.plButton.disabled = True
        popup = Popup(title="Status", content=Label(text='Workbook Closed'), size_hint=(None, None), size=(200, 100))
        popup.open()

    def restore(self):

        self.init_filters()
        self.division_drop_down.update_buttons()
        self.pl_drop_down.update_buttons()
        self.aggPkc_drop_down.update_buttons()
        self.pkgClass_drop_down.update_buttons()

        self.current_division = []
        self.current_pl = []
        self.current_pkca = []
        self.current_pkc = []

    def open_load_popup(self):
        self.upload_popup.open()

    def update_local_data(self):
        actual_quarter = prediction.utils.get_quarter()[0][-1]
        actual_year = datetime.datetime.now().year
        df = pd.read_csv(REF_PATH['ATV_All_All_All'], thousands=",")
        last_column = list(df)[-1].split("Q")
        saved_year = int(last_column[0])
        saved_quarter = int(last_column[1])

        if str(saved_quarter) != str(actual_quarter):
            print 'Updating local data...'
            if saved_quarter + 1 > 4:
                actual_quarter = 1
                actual_year = saved_year + 1
            else:
                actual_quarter = saved_quarter + 1
                actual_year = saved_year

            for key, path in REF_PATH.iteritems():
                print key,path
                df = pd.read_csv(path, thousands=",")
                header = str(actual_year)+'Q'+str(actual_quarter)

                param = key.split('_')
                success = self.macro_helper.execute_macro("change_pivot", "Division_short", [param[0]])
                self.macro_helper.execute_macro("change_pivot", "PL", [param[1]])
                self.macro_helper.execute_macro("change_pivot", "Agg_Package_PL", [param[2]])
                self.macro_helper.execute_macro("change_pivot", "PackageClass_PL", [param[3]])

                if success != 'NA':
                    xls = pd.ExcelFile(self.macro_helper.path_to_external_excel_file)
                    new_df = xls.parse('Piv_Crawl_WT_CU_pcs', skiprows=16, index_col=None, na_values=['NA']).iloc[:-1,1]
                    new_df = pd.concat([df, new_df], axis=1)
                    new_df.columns.values[-1] = header
                    if int(new_df.iloc[1,-2]) == int(new_df.iloc[1,-1]):
                        break
                    new_df.drop(new_df.columns[[1]], axis=1)
                    new_df.to_csv(path, index=False)
                else:
                    break

    def load_file(self):
        self.macro_helper.path_to_external_excel_file = self.upload_popup.file_name[0]
        print ("Loading Files...")
        self.macro_helper.open_workbook()
        self.update_local_data()
        self.init_filters()
        self.division_drop_down = DivisionDropDownList(self.division_list, self.macro_helper)
        self.aggPkc_drop_down = AggPkcDropDownList(self.pkca_list, self.macro_helper, self.mapping)
        self.pkgClass_drop_down = PkgClassDropDownList(self.pkc_list, self.macro_helper, self.mapping)
        self.pl_drop_down = PLDropDownList(self.pl_list, self.macro_helper)
        print ("Loading Complete")
        user_interface.ui.ids.divisionButton.disabled = False
        user_interface.ui.ids.pkcaButton.disabled = False
        user_interface.ui.ids.pkcButton.disabled = False
        user_interface.ui.ids.plButton.disabled = False
        popup = Popup(title="Status", content=Label(text="Loading Completed"), size_hint=(None, None), size=(200, 100))
        popup.open()

    def add_document(self):
        count = int(user_interface.ui.ids.documentCounter.text)
        user_interface.ui.ids.documentCounter.text = str(count + 1)
        download_identifiers.append(self.identifier)
        download_details.append(user_interface.ui.ids.details.text)
        download_stats.append(user_interface.ui.ids.comparisonLabel.text)
        download_images.append(self.image_path)


    @staticmethod
    def open_pop_up():
        DownloadsPopUp().open()

    @staticmethod
    def open_login_pop_up():
        LoginPopUp().open()

    def show_image_right(self):
        global window_counter
        window_counter = window_counter + 1
        if window_counter >= len(display_images):
            window_counter = 0
        user_interface.ui.ids.comparisonLabel.text = display_stats[window_counter]
        self.ids.chart.canvas.clear()
        with self.ids.chart.canvas:
            Rectangle(source=display_images[window_counter], pos=self.ids.chart.pos, size=self.ids.chart.size)

    def show_image_left(self):
        global window_counter
        window_counter = window_counter - 1
        if window_counter < 0:
            window_counter = len(display_images) - 1
        user_interface.ui.ids.comparisonLabel.text = display_stats[window_counter]
        self.ids.chart.canvas.clear()
        with self.ids.chart.canvas:
            Rectangle(source=display_images[window_counter], pos=self.ids.chart.pos, size=self.ids.chart.size)

    @staticmethod
    def get_selection(division, pl, pkca, pkc):
        if len(division) > 1:
            division = 'Multiple'
        else:
            division = division[0]
        if len(pl) > 1:
            pl = 'Multiple'
        else:
            pl = pl[0]
        if len(pkca) > 1:
            pkca = 'Multiple'
        else:
            pkca = pkca[0]
        if len(pkc) > 1:
            pkc = 'Multiple'
        else:
            pkc = pkc[0]
        return division+'_'+pl+'_'+pkca+'_'+pkc

    @staticmethod
    def get_ratio(dat, div_dat):
        return np.mean(dat.iloc[:, -4]) / np.mean(div_dat.iloc[:, -4])

    def scale_training_data(self, dat, div_dat):
        self.ratio = UI.get_ratio(dat, div_dat)
        dat.iloc[:, 1:-4] = dat.iloc[:, 1:-4] * self.ratio
        return dat

    @staticmethod
    def get_title(division, pl, pkc, pkca):
        base = 'Customer demand (potential) development pieces - OOH (CRD/WT) & sales revenue'
        base = base+'\nDIV: '+str(division)+'  PL: '+str(pl)+'  PKCA: '+str(pkca)+'  PKC: '+str(pkc)
        return base

    def submit(self):
        del display_images[:]
        del display_stats[:]
        division_lst = self.division_drop_down.current_division if len(self.division_drop_down.current_division) > 0 else ["All"]
        pl_lst = self.pl_drop_down.current_pl if len(self.pl_drop_down.current_pl) > 0 else ['All']
        pkc_lst = self.pkgClass_drop_down.current_pkc if len(self.pkgClass_drop_down.current_pkc) > 0 else ["All"]
        pkca_lst = self.aggPkc_drop_down.current_pkca if len(self.aggPkc_drop_down.current_pkca) > 0 else ["All"]

        try:
            self.get_prediction(division_lst, pl_lst, pkca_lst, pkc_lst)
            display_stats.append(user_interface.ui.ids.comparisonLabel.text)
            display_images.append(self.image_path)
            window_counter = -1
            self.show_image_right()
            popup = Popup(title="Status", content=Label(text="Prediction Completed"), size_hint=(None, None),
                          size=(200, 100))
            popup.open()
        except Exception as error:
            print error
            popup = Popup(title="Status", content=Label(text="Data Invalid"), size_hint=(None, None),
                          size=(200, 100))
            popup.open()

    def get_prediction(self, division_lst, pl_lst, pkca_lst, pkc_lst, capacity_1=None, capacity_2=None, name1="", name2=""):
        DATA_FILE = self.macro_helper.path_to_external_excel_file

        import datetime
        now = datetime.datetime.now().strftime("%Y-%m-%d %H-%M-%S")
        self.selection = UI.get_selection(division_lst, pl_lst, pkca_lst, pkc_lst)
        self.identifier = self.selection + ' ' + now
        self.image_path = PATH[:-8] + "\\output\\" + self.identifier + ".png"
        self.title = UI.get_title(division_lst,pl_lst,pkc_lst,pkca_lst)

        print ("Selected Division: ", division_lst)
        print ("Selected PL", pl_lst)
        print ("Selected PKCA: ", pkca_lst)
        print ("Selected PKC: ", pkc_lst)
        print ("Prediction ID: ", self.selection)

        if division_lst == ['All']:
            division_lst = self.division_list

        success = self.macro_helper.execute_macro("change_pivot", "Division_short", division_lst)
        self.macro_helper.execute_macro("change_pivot", "PL", pl_lst)
        self.macro_helper.execute_macro("change_pivot", "Agg_Package_PL", pkca_lst)
        self.macro_helper.execute_macro("change_pivot", "PackageClass_PL", pkc_lst)

        if success != "NA":
            all_dat = None
            all_a = None
            all_b = None
            for division in division_lst:
                if self.selection in REF_PATH.keys():
                    dat = read.data_parser.read(self.macro_helper, DATA_FILE, REF_PATH[self.selection])
                    if pkc_lst != ['All'] or pkca_lst != ['All']:
                        self.macro_helper.execute_macro("change_pivot", "Division_short", [division])
                        self.macro_helper.execute_macro("change_pivot", "PL", ['All'])
                        self.macro_helper.execute_macro("change_pivot", "Agg_Package_PL", ['All'])
                        self.macro_helper.execute_macro("change_pivot", "PackageClass_PL", ['All'])
                        div_dat = read.data_parser.read(self.macro_helper, DATA_FILE, REF_PATH[division+'_All_All_All'])
                        self.ratio = UI.get_ratio(dat, div_dat)
                    else:
                        self.ratio = 1
                else:
                    dat = read.data_parser.read(self.macro_helper, DATA_FILE, REF_PATH[division+'_All_All_All'])
                    self.macro_helper.execute_macro("change_pivot", "Division_short", [division])
                    self.macro_helper.execute_macro("change_pivot", "PL", ['All'])
                    self.macro_helper.execute_macro("change_pivot", "Agg_Package_PL", ['All'])
                    self.macro_helper.execute_macro("change_pivot", "PackageClass_PL", ['All'])
                    div_dat = read.data_parser.read(self.macro_helper, DATA_FILE, REF_PATH[division+'_All_All_All'])
                    dat = self.scale_training_data(dat, div_dat)

                a, b = prediction.predictor.predict_division(dat, division)
                if all_dat is None:
                    all_dat = dat
                    all_a = a
                    all_b = b
                else:
                    all_dat += dat
                    all_a += a
                    all_b += b

            stats = 'Prediction (next): ' + "{:,}".format(int(all_a[-1]/1000/13)) + " K/wk"
            stats += '\nPrediction (after-next): ' + "{:,}".format(int(all_b[-1]/1000/13)) + " K/wk "
            temp_act = all_dat.iloc[26,-3]
            comparison_1 = round((all_a[-1] / temp_act) * 100 - 100, 2)
            comparison_2 = round((all_b[-1] / all_a[-1]) * 100 - 100, 2)
            stats += '\nNext VS actual: ' + str(comparison_1) + '%'
            stats += '\nAfter-next VS next: ' + str(comparison_2) + '%'
            user_interface.ui.ids.comparisonLabel.text = stats

            current_input = []
            for i in range(0, len(all_b)):
                temp = []
                if i >= len(all_a):
                    temp.append(None)
                else:
                    temp.append(all_a[i])
                temp.append(all_b[i])
                current_input.append(temp)

            with open(PATH[:-8] + "\\output\\predictions.csv", "w") as my_file:
                writer = csv.writer(my_file)
                writer.writerows(current_input)

            if capacity_1 is None and capacity_2 is None:
                if user_interface.ui.ids.capacity1.text != '':

                    capacity_1 = str(user_interface.ui.ids.capacity1.text)
                    capacity_list = capacity_1.split("_")
                    if len(capacity_list) > 1:
                        capacity_1 = float(capacity_list[0])
                        name1 = "_".join(capacity_list[1:])
                    else:
                        capacity_1 = capacity_list[0]
                        name1 = ""
                else:
                    capacity_1 = 0

                if user_interface.ui.ids.capacity2.text != '':
                    capacity_2 = str(user_interface.ui.ids.capacity2.text)
                    capacity_list = capacity_2.split("_")
                    if len(capacity_list) > 1:
                        capacity_2 = float(capacity_list[0])
                        name2 = "_".join(capacity_list[1:])
                    else:
                        capacity_2 = capacity_list[0]
                        name_2 = ""
                else:
                    capacity_2 = 0

            if user_interface.ui.ids.confidence.text != '':
                level = int(user_interface.ui.ids.confidence.text)
            else:
                level = 99

            margin_next = None
            margin_afternext = None
            for div in division_lst:
                conf_path_n = PATH + "\\prediction\\param\\" + div.lower() + "_next.csv"
                conf_path_an = PATH + "\\prediction\\param\\" + div.lower() + "_afternext.csv"
                cint_n = pd.DataFrame(pd.read_csv(conf_path_n, thousands=',')).loc[:,'me'+str(level)]
                cint_an = pd.DataFrame(pd.read_csv(conf_path_an, thousands=',')).loc[:,'me'+str(level)]
                cint_n = np.append(0, cint_n)
                cint_an = np.append(0, cint_an)
                if margin_next is None:
                    margin_next = cint_n
                    margin_afternext = cint_an
                else:
                    margin_next += cint_n
                    margin_afternext += cint_an

            fig = prediction.predictor.plot_prediction(all_dat, all_a, all_b, margin_next, margin_afternext, 
                                                       ratio=self.ratio,
                                                       capa1=capacity_1, name1=name1,
                                                       capa2=capacity_2, name2=name2, title=self.title)
            plt.savefig(self.image_path)
            return stats
        else:
            return "NA"

    @staticmethod
    def clear_download():
        window_counter = -1 
        del download_identifiers[:]
        del download_details[:]
        del download_stats[:]
        del download_images[:]
        user_interface.ui.ids.documentCounter.text = "0"

    @staticmethod
    def download2():
        prs = Presentation('template.pptx')
        # Left, top, width, height
        for i in range(0, len(download_identifiers)):
            if i == 0 or i % 2 == 1:
                slide = prs.slides[-1]
            else:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
            # Details
            if i % 2 == 1:
                tx_box1 = slide.shapes.add_textbox(Inches(5.2), Inches(5.5), Inches(4.5), Inches(2))
            else:
                tx_box1 = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(4.5), Inches(2))
            tf1 = tx_box1.text_frame
            tf1.word_wrap = True
            run = tf1.paragraphs[0].add_run()
            run.text = download_details[i]
            run.font.size = Pt(14)
            # Stats
            if i % 2 == 1:
                tx_box2 = slide.shapes.add_textbox(Inches(5.2), Inches(4.2), Inches(4.5), Inches(2))
            else:
                tx_box2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(4.5), Inches(2))
            tf2 = tx_box2.text_frame
            tf2.word_wrap = True
            run = tf2.paragraphs[0].add_run()
            run.text = download_stats[i]
            run.font.size = Pt(14)
            # Graph
            if i % 2 == 1:
                slide.shapes.add_picture(download_images[i], Inches(4.8), Inches(1.2), height=Inches(3.0), width=Inches(5.0))
            else:
                slide.shapes.add_picture(download_images[i], Inches(0.1), Inches(1.2), height=Inches(3.0), width=Inches(5.0))

        date = datetime.datetime.now().strftime("%m-%d-%Y")
        prs.save(PATH[:-8] + "\\output\\" + date + '.pptx')
        del download_identifiers[:]

    def run_batch(self):
        del display_stats[:]
        del display_images[:]
        config = open(CONFIG_FILE, "r")
        keys = config.readlines()
        tasks = [str(x).strip() for x in keys]
        print (tasks)
        success = "NA"
        try:
            for task in tasks:
                print 'Predicting: '+task
                params = [i.split("_") for i in task.split(" ")]
                success = self.get_prediction(params[0], params[1], params[2], params[3], 
                                              float(params[4][0]), float(params[5][0]), 
                                              name1=params[6][0], name2=params[7][0])
                if success == "NA":
                    break
                display_images.append(self.image_path)
                display_stats.append(user_interface.ui.ids.comparisonLabel.text)
                self.add_document()
        except Exception as error:
            print error
            popup = Popup(title="Status", content=Label(text="Invalid Config"), size_hint=(None, None),
                          size=(200, 100))
            popup.open()

        if success != "NA":
            window_counter = -1
            self.show_image_right()
            UI.download2()
            self.clear_download()
            popup = Popup(title="Status", content=Label(text="Prediction Completed"), size_hint=(None, None),
                          size=(200, 100))
            popup.open()
        else:
            print "No Workbooks Opened"
            popup = Popup(title="Status", content=Label(text="No Workbooks Open"), size_hint=(None, None),
                          size=(200, 100))
            popup.open()

    drop_down = DropDown()


class ScrollableLabel(ScrollView):
    pass


class DivisionDropDownList(DropDown):
    def __init__(self, choices, macro_helper, **kwargs):
        super(DivisionDropDownList, self).__init__(**kwargs)
        self.choices = choices
        self.original = copy.deepcopy(self.choices)
        self.current_division = []
        self.macro_helper = macro_helper
        self.division_buttons = []

        self.add_items(choices)

    def add_items(self, choices):

        for choice in choices:
            btn = ToggleButton(text=choice, size_hint_y=None, height=44)
            self.add_widget(btn)
            btn.bind(on_press=self.update)
            self.division_buttons.append(btn)

        for i in self.division_buttons:
            if i in self.current_division:
                i.state = "down"

    def update_buttons(self):
        self.current_division = []
        self.clear_widgets()
        self.division_buttons = []

        self.add_items(self.choices)
        user_interface.ui.ids.divisionButton.text = "Division (" + str(len(self.choices)) + ")"

    def update(self, instance):

        if instance.state == "down":
            if instance.text not in self.current_division:
                self.current_division.append(instance.text)

        elif instance.text in self.current_division:
            self.current_division.remove(instance.text)

        count = len(self.current_division)
        if count > 1:
            user_interface.ui.ids.divisionButton.text = "Multiple Items (" + str(count) + ")"
        elif count == 1:
            user_interface.ui.ids.divisionButton.text = self.current_division[0]
        else:
            user_interface.ui.ids.divisionButton.text = "Division (" + str(len(self.choices)) + ")"


class AggPkcDropDownList(DropDown):
    def __init__(self, choices, macro_helper, mapping, **kwargs):
        super(AggPkcDropDownList, self).__init__(**kwargs)
        self.choices = choices
        self.original = copy.deepcopy(self.choices)
        self.current_pkca = []
        self.macro_helper = macro_helper
        self.pkca_buttons = []
        self.add_items(choices)
        self.mapping = mapping

    def add_items(self, choices):
        searchButton = Button(text="Click to search", size_hint_y=None, height=44, on_press=self.pkca_search_pop)
        self.add_widget(searchButton)

        for choice in choices:
            btn = ToggleButton(text=choice, size_hint_y=None, height=44)
            self.add_widget(btn)
            btn.bind(on_press=self.update)
            self.pkca_buttons.append(btn)

        for i in self.pkca_buttons:
            if i in self.current_pkca:
                i.state = "down"
            
    def pkca_search_pop(self, instance):
        self.popup = Popup(title="PKCA Search", size_hint = (0.8,0.7),
                           content=TextInput(hint_text = "type search here; to get back all choices, type restore ", multiline = False, 
                           on_text_validate= self.pkca_search, size_hint=(0.7,0.3)))
        self.popup.open()



    def pkca_search(self, instance):
        text = self.popup.content.text

        if text == "restore" or text == '' or text == ' ':
            self.update_buttons()

        else:
            temp = []
            for choice in self.choices:
                if text in choice:
                    temp.append(choice)
            self.update_buttons(choices=temp)



    def update_buttons(self, choices=None):

        if choices is None:
            self.current_pkca = []
            self.clear_widgets()
            self.pkca_buttons = []
            self.add_items(self.choices)
            user_interface.ui.ids.pkcaButton.text = "PKCA (" + str(len(self.choices)) + ")"
        else:
            if len(self.current_pkca) == 0:
                self.current_pkca = []
                self.clear_widgets()
                self.pkca_buttons = []
                self.add_items(choices)
                user_interface.ui.ids.pkcaButton.text = "PKCA (" + str(len(self.choices)) + ")"

    def update(self, instance):
        if instance.state == "down":
            if instance.text not in self.current_pkca:
                self.current_pkca.append(instance.text)
        elif instance.text in self.current_pkca:
            self.current_pkca.remove(instance.text)

        count = len(self.current_pkca)
        if count > 1:
            user_interface.ui.ids.pkcaButton.text = "Multiple Items (" + str(count) + ")"
        elif count == 1:
            user_interface.ui.ids.pkcaButton.text = self.current_pkca[0]
        else:
            user_interface.ui.ids.pkcaButton.text = "PKCA (" + str(len(self.choices)) + ")"

        if len(self.current_pkca) == 0:
            user_interface.ui.pkgClass_drop_down.update_buttons(choices=None)
        else:
            new_mapping = self.mapping.loc[self.mapping["PKCA"] == self.current_pkca[0]]
            valid_options = list(set(list(new_mapping["PKC"])))
            valid_options.sort()
            user_interface.ui.pkgClass_drop_down.update_buttons(choices=valid_options)


class PkgClassDropDownList(DropDown):
    def __init__(self, choices, macro_helper, mapping, **kwargs):
        super(PkgClassDropDownList, self).__init__(**kwargs)
        self.mapping = mapping
        self.choices = choices
        self.original = copy.deepcopy(self.choices)
        self.current_pkc = []
        self.macro_helper = macro_helper
        self.pkc_buttons = []
        self.add_items(choices)

    def add_items(self, choices):
        searchButton = Button(text="Click to search", size_hint_y=None, height=44, on_press=self.pkc_search_pop)
        self.add_widget(searchButton)
	
        for choice in choices:
            btn = ToggleButton(text=choice, size_hint_y=None, height=44)
            self.add_widget(btn)
            btn.bind(on_press=self.update)
            self.pkc_buttons.append(btn)

        for i in self.pkc_buttons:
            if i in self.current_pkc:
                i.state = "down"

    def pkc_search_pop(self, instance):
        self.popup = Popup(title="PKC Search", size_hint = (0.8,0.7),
                           content=TextInput(hint_text = "type search here; to get back all choices, type restore", multiline = False, 
                           on_text_validate= self.pkc_search, size_hint=(0.7,0.3)))
        self.popup.open()

    def pkc_search(self, instance):
        text = self.popup.content.text

        if text == "restore" or text == '' or text == ' ':
            self.update_buttons()
        else:
            temp = []
            for choice in self.choices:
                if text in choice:
                    temp.append(choice)
            self.update_buttons(choices=temp)


    def update_buttons(self, choices=None):
        if choices is None:
            self.current_pkc = []
            self.clear_widgets()
            self.pkc_buttons = []
            self.add_items(self.choices)
            user_interface.ui.ids.pkcButton.text = "PKC (" + str(len(self.choices)) + ")"
        else:
            if len(self.current_pkc) == 0:
                self.current_pkc = []
                self.clear_widgets()
                self.pkc_buttons = []
                self.add_items(choices)
                user_interface.ui.ids.pkcButton.text = "PKC (" + str(len(self.choices)) + ")"

    def update(self, instance):
        if instance.state == "down":
            if instance.text not in self.current_pkc:
                self.current_pkc.append(instance.text)
        elif instance.text in self.current_pkc:
            self.current_pkc.remove(instance.text)

        count = len(self.current_pkc)
        if count > 1:
            user_interface.ui.ids.pkcButton.text = "Multiple Items (" + str(count) + ")"
        elif count == 1:
            user_interface.ui.ids.pkcButton.text = self.current_pkc[0]
        else:
            user_interface.ui.ids.pkcButton.text = "PKC (" + str(len(self.choices)) + ")"

        if len(self.current_pkc) == 0:
            user_interface.ui.aggPkc_drop_down.update_buttons(choices=None)
        else:
            new_mapping = self.mapping.loc[self.mapping["PKC"] == self.current_pkc[0]]
            valid_options = list(set(list(new_mapping["PKCA"])))
            valid_options.sort()
            user_interface.ui.aggPkc_drop_down.update_buttons(choices=valid_options)


class PLDropDownList(DropDown):
    def __init__(self, choices, macro_helper, **kwargs):
        super(PLDropDownList, self).__init__(**kwargs)
        self.choices = choices
        self.original = copy.deepcopy(self.choices)
        self.current_pl = []
        self.macro_helper = macro_helper
        self.pl_buttons = []

        self.add_items(choices)

    def add_items(self, choices):
        searchButton = Button(text="Click to search", size_hint_y=None, height=44, on_press=self.pl_search_pop)
        self.add_widget(searchButton)

        for choice in choices:
            btn = ToggleButton(text=choice, size_hint_y=None, height=44)
            self.add_widget(btn)
            btn.bind(on_press=self.update)
            self.pl_buttons.append(btn)

        for i in self.pl_buttons:
            if i in self.current_pl:
                i.state = "down"

    def pl_search_pop(self, instance):
        self.popup = Popup(title="PL Search", size_hint = (0.8,0.7),
                           content=TextInput(hint_text = "type search here; to get back all choices, type restore", multiline = False,
                           on_text_validate= self.pl_search, size_hint=(0.7,0.3)))
        self.popup.open()

    def pl_search(self, instance):
        text = self.popup.content.text

        if text == "restore" or text == '' or text == ' ':
            self.update_buttons()
        else:
            temp = []
            for choice in self.choices:
                if text in choice:
                    temp.append(choice)
            self.update_buttons(choices=temp)

    def update_buttons(self, choices=None):
        if choices is None:
            self.current_pl = []
            self.clear_widgets()
            self.pl_buttons = []
            self.add_items(self.choices)
            user_interface.ui.ids.plButton.text = "PL (" + str(len(self.choices)) + ")"
        else:
            if len(self.current_pl) == 0:
                self.current_pl = []
                self.clear_widgets()
                self.pl_buttons = []
                self.add_items(choices)
                user_interface.ui.ids.plButton.text = "PL (" + str(len(self.choices)) + ")"

    def update(self, instance):
        if instance.state == "down":
            if instance.text not in self.current_pl:
                self.current_pl.append(instance.text)
        elif instance.text in self.current_pl:
            self.current_pl.remove(instance.text)

        count = len(self.current_pl)
        if count > 1:
            user_interface.ui.ids.plButton.text = "Multiple Items (" + str(count) + ")"
        elif count == 1:
            user_interface.ui.ids.plButton.text = self.current_pl[0]
        else:
            user_interface.ui.ids.plButton.text = "PL (" + str(len(self.choices)) + ")"


class Downloads(BoxLayout):
    def __init__(self, **kwargs):
        super(Downloads, self).__init__(**kwargs)
        user_interface.ui.ids.documentCounter.text = str(len(download_identifiers))
        del download_buttons[:]
        if len(download_identifiers) != 0:
            self.layout = BoxLayout(orientation='vertical')
            for i in range(0, len(download_identifiers)):
                layout1 = BoxLayout()
                label = Label(text=download_identifiers[i])
                button = Button(text="remove")
                button.bind(on_press=self.remove)
                download_buttons.append(button)
                layout1.add_widget(label)
                layout1.add_widget(button)
                self.layout.add_widget(layout1)

            download_btn = Button(text="download")
            download_btn.bind(on_press=self.download)
            self.layout.add_widget(download_btn)
            self.add_widget(self.layout)

    def download(self, instance):
        prs = Presentation('template.pptx')
        # Left, top, width, height
        for i in range(0, len(download_identifiers)):
            if i == 0 or i % 2 == 1:
                slide = prs.slides[-1]
            else:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
            # Details
            if i % 2 == 1:
                tx_box1 = slide.shapes.add_textbox(Inches(5.2), Inches(5.5), Inches(4.5), Inches(2))
            else:
                tx_box1 = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(4.5), Inches(2))
            tf1 = tx_box1.text_frame
            tf1.word_wrap = True
            run = tf1.paragraphs[0].add_run()
            run.text = download_details[i]
            run.font.size = Pt(14)
            # Stats
            if i % 2 == 1:
                tx_box2 = slide.shapes.add_textbox(Inches(5.2), Inches(4.2), Inches(4.5), Inches(2))
            else:
                tx_box2 = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(4.5), Inches(2))
            tf2 = tx_box2.text_frame
            tf2.word_wrap = True
            run = tf2.paragraphs[0].add_run()
            run.text = download_stats[i]
            run.font.size = Pt(14)
            # Graph
            if i % 2 == 1:
                slide.shapes.add_picture(download_images[i], Inches(4.8), Inches(1.2), height=Inches(3.0), width=Inches(5.0))
            else:
                slide.shapes.add_picture(download_images[i], Inches(0.1), Inches(1.2), height=Inches(3.0), width=Inches(5.0))

        date = datetime.datetime.now().strftime("%m-%d-%Y")
        prs.save(PATH[:-8] + "\\output\\" + date + '.pptx')
        del download_identifiers[:]

        del download_images[:]
        del download_details[:]
        del download_stats[:]
        self.update()

    # update the downloads shown
    def update(self):
        self.layout.clear_widgets()
        user_interface.ui.ids.documentCounter.text = str(len(download_identifiers))
        if len(download_identifiers) != 0:
            for i in range(0, len(download_identifiers)):
                layout1 = BoxLayout()
                label = Label(text=download_identifiers[i])
                button = Button(text="remove")
                button.bind(on_press=self.remove)
                download_buttons.append(button)
                layout1.add_widget(label)
                layout1.add_widget(button)
                self.layout.add_widget(layout1)

            download_btn = Button(text="download")
            download_btn.bind(on_press=self.download)
            self.layout.add_widget(download_btn)

    # remove the specific download shown
    def remove(self, instance):
        for i in range(0, len(download_identifiers)):
            if download_buttons[i] == instance:
                del download_identifiers[i]
                del download_details[i]
                del download_stats[i]
                del download_images[i]
                del download_buttons[:]
                break
        self.update()


class DownloadsPopUp(Popup):
    def show(self):
        self.content = Downloads()
        return self.content



# upload file popup
class UploadPopUp(Popup):

    def __init__(self, **kwargs):
        super(UploadPopUp, self).__init__(**kwargs)
        self.file_path = None
        self.file_name = None

    def cancel(self):
        self.file_path = None
        self.file_name = None
        self.dismiss()

    def load(self, file_path, file_name):
        self.file_path = file_path
        self.file_name = file_name
        user_interface.ui.load_file()
        self.dismiss()


class UserInterface(App):
    kv_directory = 'ui'

    def __init__(self, **kwargs):
        super(UserInterface, self).__init__(**kwargs)
        self.ui = UI()

    def build(self):
        return self.ui


if __name__ == "__main__":
    user_interface = UserInterface()
    user_interface.run()
